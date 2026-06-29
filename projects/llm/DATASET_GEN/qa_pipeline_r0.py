#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Q&A dataset generator with Smart OCR support for PDF, DOC/DOCX, and PPT/PPTX.

This version is configured to run **entirely offline** using the EMTAC .env file:
    E:\emtac\dev_env\.env

Environment variables control:
    - Hugging Face model cache paths
    - Tesseract OCR path and tessdata directory
    - Poppler path for PDF OCR
"""

import os
import re
import sys
import json
import time
import math
import argparse
import logging
import datetime
from io import BytesIO
from pathlib import Path
from typing import List, Dict, Optional, Tuple
import dotenv
# -----------------------------------------
# Logging
# -----------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)-7s | %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S"
)
log = logging.getLogger("qa_pipeline")

# -----------------------------------------
# Load .env first (CRITICAL)
# -----------------------------------------
from dotenv import load_dotenv

ENV_PATH = r"E:\emtac\dev_env\.env"
if os.path.exists(ENV_PATH):
    load_dotenv(dotenv_path=ENV_PATH, override=True)
    log.info(f"[ENV] Loaded environment configuration from: {ENV_PATH}")
else:
    log.warning(f"[ENV] File not found at {ENV_PATH}. Using process environment variables.")

# -----------------------------------------
# Third-party dependencies
# -----------------------------------------
import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pdfplumber
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
import torch
import pandas as pd
from tqdm.auto import tqdm
from sentence_transformers import SentenceTransformer, util
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
import re

from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from openpyxl.styles import Alignment, Font, PatternFill


SENTENCE_SPLIT_RE = re.compile(r"(?<=[.?!])\s+(?=[A-Z])")
# -----------------------------------------
# Offline / Environment Settings
# -----------------------------------------
os.environ["TRANSFORMERS_OFFLINE"] = "1"
os.environ["HF_DATASETS_OFFLINE"] = "1"
os.environ["HF_HOME"] = os.environ.get("HF_HOME", r"E:\emtac\models")
os.environ["HF_HUB_CACHE"] = os.environ.get("HF_HUB_CACHE", r"E:\emtac\models\cache")
os.environ["TRANSFORMERS_CACHE"] = os.environ.get("TRANSFORMERS_CACHE", r"E:\emtac\models\cache")

# -----------------------------------------
# OCR Path Configuration
# -----------------------------------------
TESSERACT_CMD = os.environ.get("TESSERACT_CMD")
if TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
else:
    log.error("TESSERACT_CMD not set in .env (OCR will fail).")

TESSDATA_PREFIX = os.environ.get("TESSDATA_PREFIX")
if TESSDATA_PREFIX:
    os.environ["TESSDATA_PREFIX"] = TESSDATA_PREFIX
else:
    log.warning("TESSDATA_PREFIX not set in .env. OCR for languages beyond English may fail.")

POPPLER_BIN = os.environ.get("POPPLER_BIN")
if not POPPLER_BIN:
    log.error("POPPLER_BIN not set in .env — PDF OCR will fail.")

log.info(f"[OCR] Tesseract Executable: {TESSERACT_CMD}")
log.info(f"[OCR] Tessdata Directory:   {TESSDATA_PREFIX}")
log.info(f"[OCR] Poppler Bin Folder:   {POPPLER_BIN}")

# -----------------------------------------
# Local Model Paths
# -----------------------------------------
MINILM_LOCAL = os.environ.get("MODEL_MINILM_DIR", r"E:\emtac\models\cache\all-MiniLM-L6-v2")
FLAN_LOCAL = os.environ.get("MODEL_LLM_DIR", r"E:\emtac\models\llm\flan_t5_large\models--google--flan-t5-large\snapshots\0613663d0d48ea86ba8cb3d7a44f0f65dc596a2a")


# -----------------------------------------
# Configuration
# -----------------------------------------
HEADING_RE = re.compile(
    r'^(?:[A-Z][A-Z0-9 ._-]{3,}|[0-9]+[.)][^\n]{3,})$',
    re.MULTILINE
)

CHUNK_SIZE = 120
SMART_OCR_CHAR_THRESHOLD = 200
OCR_LANG = "eng"
TEST_FILE = r"E:\emtac\data\raw_documention\FB4-GENERAL\BOL Checklist Part Fill.pdf"

# -----------------------------------------
# OCR Helpers
# -----------------------------------------

def ocr_image_bytes(img_bytes: bytes, lang: str = OCR_LANG) -> str:
    with Image.open(BytesIO(img_bytes)) as im:
        im = im.convert("RGB")
        return pytesseract.image_to_string(im, lang=lang)

def ocr_pil_image(img: Image.Image, lang: str = OCR_LANG) -> str:
    im = img.convert("RGB")
    return pytesseract.image_to_string(im, lang=lang)

def clean_text(s: str) -> str:
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\r\n", "\n", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


# -----------------------------------------
# Pipeline Class
# -----------------------------------------

class QAPipeline:
    def __init__(self, device=None, ocr_lang=OCR_LANG):
        self.device = device or ("cuda" if torch.cuda.is_available() else "cpu")
        self.ocr_lang = ocr_lang

        log.info("Loading FLAN-T5 (offline)...")
        self.answer_tokenizer = AutoTokenizer.from_pretrained(FLAN_LOCAL, local_files_only=True)
        self.answer_model = AutoModelForSeq2SeqLM.from_pretrained(FLAN_LOCAL, local_files_only=True).to(self.device)

        log.info("Loading MiniLM Embedding (offline)...")
        self.embed_model = SentenceTransformer(MINILM_LOCAL, device=self.device)

        # --- T5 Question Generation model
        T5_QG_DIR = os.environ.get("MODEL_mrm8488_DIR")
        if not T5_QG_DIR or not os.path.isdir(T5_QG_DIR):
            raise RuntimeError("MODEL_mrm8488_DIR is not set or invalid. Set it to your local T5 QG path.")
        log.info(f"Loading T5 QG from: {T5_QG_DIR}")
        self.qg_tokenizer = AutoTokenizer.from_pretrained(T5_QG_DIR, local_files_only=True)
        self.qg_model = AutoModelForSeq2SeqLM.from_pretrained(T5_QG_DIR, local_files_only=True).to(self.device)

    # -------------------------------------
    # Document Loaders (Smart OCR)
    # -------------------------------------
    def load_txt(self, path: str) -> str:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return clean_text(f.read())

    def load_docx(self, path: str) -> str:
        doc = docx.Document(path)
        base_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        base_text = clean_text(base_text)

        images_bytes = [
            rel._target.part.blob
            for rel in doc.part.rels.values()
            if "image" in rel.reltype
        ]

        if len(base_text) >= SMART_OCR_CHAR_THRESHOLD:
            return base_text

        ocr_blocks = [ocr_image_bytes(b) for b in images_bytes]
        return clean_text(base_text + "\n\n" + "\n\n".join(ocr_blocks))

    def load_pptx(self, path: str) -> str:
        pres = Presentation(path)
        slides = []
        for slide in pres.slides:
            text_parts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            slide_text = clean_text("\n".join(text_parts))
            if len(slide_text) < SMART_OCR_CHAR_THRESHOLD:
                ocr_blocks = []
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        ocr_blocks.append(ocr_image_bytes(shape.image.blob))
                slide_text += "\n\n" + "\n\n".join(ocr_blocks)
            slides.append(slide_text)
        return clean_text("\n\n".join(slides))

    def load_pdf(self, path: str) -> str:
        pages_out = []
        with pdfplumber.open(path) as pdf:
            for pg, page in enumerate(pdf.pages, start=1):
                txt = clean_text(page.extract_text() or "")
                if len(txt) >= SMART_OCR_CHAR_THRESHOLD:
                    pages_out.append(txt)
                    continue
                imgs = convert_from_path(path, dpi=300, first_page=pg, last_page=pg, poppler_path=POPPLER_BIN)
                pages_out.append(clean_text(txt + "\n\n" + ocr_pil_image(imgs[0], self.ocr_lang)))
        return clean_text("\n\n".join(pages_out))

    # -------------------------------------
    # Text Chunking, Q&A, Dedup, Output
    # -------------------------------------
    def chunk_text(self, text: str, overlap: int = 0) -> List[str]:
        words = text.split()
        chunks = []
        step = CHUNK_SIZE if overlap == 0 else max(1, CHUNK_SIZE - overlap)
        for i in range(0, len(words), step):
            chunk = " ".join(words[i:i + CHUNK_SIZE])
            if len(chunk.split()) > 20:
                chunks.append(chunk)
        return chunks

    @torch.no_grad()
    def generate_questions_t5(self, context: str, max_questions: int = 5, max_answer_words: int = 30) -> List[
        Tuple[str, str]]:
        """
        Treat each sentence in the context as a candidate 'answer'.
        For each, ask T5 to produce a question conditioned on (answer, context).
        Returns (question, answer) pairs. Truncates overly long answers.
        """
        pairs: List[Tuple[str, str]] = []
        sentences = self.split_sentences(context)

        # Light prioritization: earlier sentences first
        for sent in sentences:
            # Trim very long answers to keep the pair crisp
            awords = sent.split()
            if len(awords) > max_answer_words:
                sent = " ".join(awords[:max_answer_words])

            prompt = f"answer: {sent}  context: {context}"
            inputs = self.qg_tokenizer(
                prompt, return_tensors="pt",
                truncation=True, max_length=512
            ).to(self.device)

            outputs = self.qg_model.generate(
                **inputs,
                num_beams=4,
                max_new_tokens=64,
                early_stopping=True
            )
            question = self.qg_tokenizer.decode(outputs[0], skip_special_tokens=True).strip()

            # --- Clean common prefix artifacts ---
            # remove "question:" or "Question:" if present
            question = re.sub(r"^(question\s*:)\s*", "", question, flags=re.IGNORECASE).strip()

            # ensure trailing ?
            if not question.endswith("?"):
                question += "?"

            # Basic quality filters
            if len(question.split()) < 3:
                continue
            if not question.endswith("?"):
                question += "?"

            pairs.append((question, sent))
            if len(pairs) >= max_questions:
                break

        return pairs

    def trim_answer_phrase(self, answer: str) -> str:
        answer = (answer or "").strip()
        if not answer:
            return answer

        # Preserve conditional openings like "If ..."
        if answer.lower().startswith("if "):
            return answer.rstrip(" .")

        # Remove one leading subject phrase (e.g., "Operators", "The conveyor")
        answer = re.sub(r"^(?:[A-Z][a-z]+|The|the)\s+(?:[A-Z]?[a-z]+)\s+", "", answer, count=1)

        # Trim and clean
        answer = answer.strip().rstrip(" .")

        # Lowercase initial char if not starting with "If"
        if answer and not answer.lower().startswith("if "):
            answer = answer[0].lower() + answer[1:]

        return answer

    @torch.no_grad()
    def clean_answer_with_flan(self, question: str, rough_answer: str, context: str) -> str:
        """
        Re-aligns answer so that it exactly matches a span of the context.
        This prevents truncation or hallucinated filler words.
        """
        prompt = (
            "You are a precision extraction model.\n"
            "Select the exact phrase from the context that answers the question.\n"
            "Return text copied *only* from the context.\n\n"
            f"CONTEXT:\n{context}\n\n"
            f"QUESTION: {question}\n"
            f"ROUGH_ANSWER: {rough_answer}\n\n"
            "Return the corrected answer:"
        )

        inputs = self.answer_tokenizer(prompt, return_tensors="pt", truncation=True).to(self.device)
        outputs = self.answer_model.generate(inputs["input_ids"], max_new_tokens=40)
        cleaned = self.answer_tokenizer.decode(outputs[0], skip_special_tokens=True).strip()

        # If cleaned answer actually comes from context, use it
        if cleaned.lower() in context.lower():
            return cleaned

        # Otherwise, keep the rough version
        return rough_answer

    def deduplicate(self, qa_pairs: List[Dict]) -> List[Dict]:
        if not qa_pairs:
            return []
        questions = [p["question"] for p in qa_pairs]
        embeddings = self.embed_model.encode(questions, convert_to_tensor=True)
        keep, used = [], set()
        for i in range(len(qa_pairs)):
            if i in used:
                continue
            keep.append(qa_pairs[i])
            scores = util.pytorch_cos_sim(embeddings[i], embeddings)[0]
            dupes = torch.where(scores > 0.90)[0]
            for d in dupes:
                used.add(int(d))
        return keep

    def save_outputs(self, qa_pairs: List[Dict], out_dir: str, base_name: str, src_path: str, timestamp: str):
        os.makedirs(out_dir, exist_ok=True)

        # JSONL
        json_path = os.path.join(out_dir, f"{base_name}.jsonl")
        with open(json_path, "w", encoding="utf-8") as f:
            for item in qa_pairs:
                f.write(json.dumps(item, ensure_ascii=False) + "\n")

        # Excel
        df = pd.DataFrame(qa_pairs)
        xlsx_path = os.path.join(out_dir, f"{base_name}.xlsx")
        df.to_excel(xlsx_path, index=False)

        # Text file
        txt_path = os.path.join(out_dir, f"{base_name}.txt")
        self.save_as_text(qa_pairs, txt_path, src_path, timestamp)

        log.info(f"Saved JSONL → {json_path}")
        log.info(f"Saved Excel → {xlsx_path}")
        log.info(f"Saved TXT → {txt_path}")

    def save_as_text(self, qa_pairs: List[Dict], txt_path: str, src_path: str, timestamp: str):
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write("Q&A DATASET SUMMARY\n")
            f.write("-------------------\n")
            f.write(f"Source Document: {src_path}\n")
            f.write(f"Generated On:    {timestamp}\n")
            f.write(f"Total Q&A Pairs: {len(qa_pairs)}\n\n")

            f.write("Q&A SET (Readable Format)\n")
            f.write("-------------------------\n\n")

            for qa in qa_pairs:
                q = (qa.get("question") or "").strip()
                ans_final = (qa.get("answer") or "").strip()
                ans_rough = (qa.get("answer_rough") or "").strip()  # <-- fixed key
                ctx = (qa.get("context") or "").strip()

                f.write(f"Q: {q}\n")
                if ans_rough and ans_rough != ans_final:
                    f.write(f"A (Rough): {ans_rough}\n")
                    f.write(f"A (Final): {ans_final}\n")
                else:
                    f.write(f"A: {ans_final}\n")

                # Keep context visible for traceability
                f.write(f"Context: {ctx}\n\n")

    def build(self, src_path: str, output_path: Optional[str] = None):
        src_path = str(src_path)
        ext = Path(src_path).suffix.lower()

        # --- Load document according to type ---
        if ext == ".docx":
            text = self.load_docx(src_path)
        elif ext in (".pptx", ".ppt"):
            text = self.load_pptx(src_path)
        elif ext == ".pdf":
            text = self.load_pdf(src_path)
        elif ext == ".txt":
            text = self.load_txt(src_path)
        else:
            raise ValueError(f"Unsupported type: {ext}")

        # --- Create overlapping chunks to improve Q coverage ---
        pass1 = self.chunk_text(text, overlap=0)
        pass2 = self.chunk_text(text, overlap=60)
        chunks = pass1 + pass2

        qa_pairs = []

        # --- MAIN QG LOOP (T5-based) ---
        for chunk in tqdm(chunks, desc="Generating Q&A..."):
            # generate up to 4 (question, answer) pairs per chunk
            for q, a in self.generate_questions_t5(chunk, max_questions=4):
                if a.lower() in chunk.lower():
                    # Rough answer
                    rough = self.trim_answer_phrase(a)

                    # Smoothed / corrected answer
                    smooth = self.clean_answer_with_flan(q, rough, chunk)

                    qa_pairs.append({
                        "question": q,
                        "answer_rough": rough,
                        "answer": smooth,
                        "context": chunk
                    })

        # --- Deduplicate near-duplicate questions ---
        qa_pairs = self.deduplicate(qa_pairs)

        # --- Save output ---
        base_name = Path(src_path).stem
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        out_dir = output_path or os.path.join(base_name, timestamp)
        self.save_outputs(qa_pairs, out_dir, base_name, src_path, timestamp)
        log.info(f"Done. Output directory:\n{out_dir}")

    def split_sentences(self, text: str) -> List[str]:
        raw = re.split(SENTENCE_SPLIT_RE, text)
        return [s.strip() for s in raw if len(s.strip().split()) >= 6]
# -----------------------------------------
# CLI Entrypoint
# -----------------------------------------

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("source", nargs="?", help="Document to process")
    parser.add_argument("--out", help="Output directory override", default=None)
    parser.add_argument("--test", action="store_true",
                        help=f"Run on test fixture: {TEST_FILE}")
    parser.add_argument("--check", action="store_true", help="Run system diagnostic and exit")
    args = parser.parse_args()

    if args.check:
        system_check()
        return

    # --- TEST MODE ---
    if args.test:
        src_path = TEST_FILE
        if not os.path.exists(src_path):
            print(f"[ERROR] Test file not found:\n{src_path}")
            return
        print(f"[TEST MODE] Using test document: {src_path}")
        pipe = QAPipeline()
        pipe.build(src_path, output_path=args.out)
        return

    # --- NORMAL MODE ---
    if not args.source:
        print("Error: You must provide a source file unless using --test or --check")
        return

    pipe = QAPipeline()
    pipe.build(args.source, output_path=args.out)


def system_check():
    print("\n=== EMTAC AI SYSTEM DIAGNOSTIC CHECK ===\n")

    # --- .env paths ---
    print("[ENVIRONMENT VARIABLES]")
    keys = [
        "HF_HOME", "HF_HUB_CACHE", "TRANSFORMERS_CACHE",
        "TESSERACT_CMD", "TESSDATA_PREFIX", "POPPLER_BIN",
        "MODEL_LLM_DIR", "MODEL_MINILM_DIR"
    ]
    for k in keys:
        print(f"  {k}: {os.environ.get(k, 'NOT SET')}")

    print("\n[FILE PATH VALIDATION]")

    def exists(path, name):
        if path and os.path.exists(path):
            print(f"  ✔ {name}: {path}")
        else:
            print(f"  ✖ {name}: NOT FOUND → {path}")

    exists(os.environ.get("TESSERACT_CMD"), "Tesseract EXE")
    exists(os.environ.get("TESSDATA_PREFIX"), "Tessdata Folder")
    exists(os.environ.get("POPPLER_BIN"), "Poppler Bin Directory")
    exists(os.environ.get("MODEL_LLM_DIR"), "FLAN Model Directory")
    exists(os.environ.get("MODEL_MINILM_DIR"), "MiniLM Embedding Model Directory")

    print("\n[OCR FUNCTIONAL TEST]")
    try:
        cmd = os.environ.get("TESSERACT_CMD")
        out = os.popen(f'"{cmd}" --version').read()
        if "tesseract" in out.lower():
            print("  ✔ Tesseract is working")
        else:
            print("  ✖ Tesseract failed to execute")
    except Exception as e:
        print(f"  ✖ Error running tesseract: {e}")

    tessdata = os.environ.get("TESSDATA_PREFIX")
    if tessdata:
        langfile = Path(tessdata) / "eng.traineddata"
        print(f"  {'✔' if langfile.exists() else '✖'} Tessdata contains eng.traineddata")

    print("\n[POPPLER FUNCTIONAL TEST]")
    try:
        test_pdf = r"E:\emtac\data\test_check.pdf"
        if Path(test_pdf).exists():
            pages = convert_from_path(test_pdf, dpi=50, first_page=1, last_page=1, poppler_path=os.environ.get("POPPLER_BIN"))
            print(f"  ✔ Poppler converted 1 test page")
        else:
            print("  (Skipping - place a PDF at E:\\emtac\\data\\test_check.pdf to verify)")
    except Exception as e:
        print(f"  ✖ Poppler conversion failed: {e}")

    print("\n[MODEL LOAD TEST]")
    try:
        _tok = AutoTokenizer.from_pretrained(os.environ.get("MODEL_LLM_DIR"), local_files_only=True)
        _mdl = AutoModelForSeq2SeqLM.from_pretrained(os.environ.get("MODEL_LLM_DIR"), local_files_only=True)
        print("  ✔ FLAN-T5 model loaded offline successfully")
    except Exception as e:
        print(f"  ✖ FLAN model load FAILED: {e}")

    try:
        _embed = SentenceTransformer(os.environ.get("MODEL_MINILM_DIR"))
        print("  ✔ MiniLM embedding model loaded offline successfully")
    except Exception as e:
        print(f"  ✖ MiniLM load FAILED: {e}")

    print("\n=== SYSTEM CHECK COMPLETE ===\n")


if __name__ == "__main__":
    main()

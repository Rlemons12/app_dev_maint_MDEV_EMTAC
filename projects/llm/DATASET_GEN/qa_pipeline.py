#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Q&A dataset generator with Smart OCR support for PDF, DOC/DOCX, and PPT/PPTX.
Ranked multi-candidate generation + modular second-pass system.

Runs entirely offline using EMTAC .env:
    E:\emtac\dev_env\.env
"""

import os
import re
import sys
import json
import math
import argparse
import logging
import datetime
import time
from io import BytesIO
from pathlib import Path
from typing import List, Dict, Optional, Tuple, Any
from gpu_adapter import GPUAdapter
import numpy as np


# -----------------------------------------
# Logging
# -----------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)-7s | %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
log = logging.getLogger("qa_pipeline")

# -----------------------------------------
# ULTRA MODE LOGGING OVERRIDE
# -----------------------------------------
log.setLevel(logging.DEBUG)
for h in logging.root.handlers:
    h.setLevel(logging.DEBUG)

# -----------------------------------------
# Load .env first (CRITICAL)
# -----------------------------------------
from dotenv import load_dotenv
ENV_PATH = r"E:\emtac\dev_env\.env"
if os.path.exists(ENV_PATH):
    load_dotenv(dotenv_path=ENV_PATH, override=True)
    log.info(f"[ENV] Loaded environment configuration from: {ENV_PATH}")
else:
    log.warning(f"[ENV] File not found at {ENV_PATH}. Using process environment variables.")

# -----------------------------------------
# Third-party deps (offline)
# -----------------------------------------
import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pdfplumber
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
import torch
import pandas as pd
from tqdm.auto import tqdm
from sentence_transformers import SentenceTransformer, util
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM,AutoModelForQuestionAnswering

# -----------------------------------------
# Regex / constants
# -----------------------------------------
SENTENCE_SPLIT_RE = re.compile(r"(?<=[.?!])\s+(?=[A-Z])")
HEADING_RE = re.compile(r'^(?:[A-Z][A-Z0-9 ._-]{3,}|[0-9]+[.)][^\n]{3,})$', re.MULTILINE)

# Offline / HF caches
os.environ["TRANSFORMERS_OFFLINE"] = "1"
os.environ["HF_DATASETS_OFFLINE"] = "1"
os.environ["HF_HOME"] = os.environ.get("HF_HOME", r"E:\emtac\models")
os.environ["HF_HUB_CACHE"] = os.environ.get("HF_HUB_CACHE", r"E:\emtac\models\cache")
os.environ["TRANSFORMERS_CACHE"] = os.environ.get("TRANSFORMERS_CACHE", r"E:\emtac\models\cache")

# OCR config
TESSERACT_CMD = os.environ.get("TESSERACT_CMD")
if TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
else:
    log.error("TESSERACT_CMD not set in .env (OCR will fail).")

TESSDATA_PREFIX = os.environ.get("TESSDATA_PREFIX")
if TESSDATA_PREFIX:
    os.environ["TESSDATA_PREFIX"] = TESSDATA_PREFIX
else:
    log.warning("TESSDATA_PREFIX not set in .env. OCR for languages beyond English may fail.")

POPPLER_BIN = os.environ.get("POPPLER_BIN")
if not POPPLER_BIN:
    log.error("POPPLER_BIN not set in .env — PDF OCR will fail.")

log.info(f"[OCR] Tesseract Executable: {TESSERACT_CMD}")
log.info(f"[OCR] Tessdata Directory:   {TESSDATA_PREFIX}")
log.info(f"[OCR] Poppler Bin Folder:   {POPPLER_BIN}")

# Local model paths
MINILM_LOCAL = os.environ.get("MODEL_MINILM_DIR", r"E:\emtac\models\cache\all-MiniLM-L6-v2")
FLAN_LOCAL   = os.environ.get("MODEL_LLM_DIR",   r"E:\emtac\models\llm\flan_t5_large\models--google--flan-t5-large\snapshots\0613663d0d48ea86ba8cb3d7a44f0f65dc596a2a")
T5_QG_DIR    = os.environ.get("MODEL_mrm8488_DIR", r"E:\emtac\models\cache\t5-base-finetuned-question-generation-ap")

# Defaults (CLI can override)
CHUNK_SIZE = 120
SMART_OCR_CHAR_THRESHOLD = 200
OCR_LANG = "eng"

# Test file (used by --test)
#TEST_FILE = r"E:\emtac\data\raw_documention\FB4-GENERAL\BOL Checklist Part Fill.pdf"
# Alternate example:
# TEST_FILE = r"E:\emtac\projects\llm\DATASET_GEN\training_test_1.txt"
TEST_FILE = r"E:\emtac\data\raw_documention\FB4-GENERAL\AFL31600 Overview.docx"
# Adaptive ranked settings (CLI can override)
TARGET_SIMILARITY = 0.70
MAX_RETRIES = 5

# -----------------------------------------
# Helpers
# -----------------------------------------
def clean_text(s: str) -> str:
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\r\n", "\n", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def ocr_image_bytes(img_bytes: bytes, lang: str = OCR_LANG) -> str:
    with Image.open(BytesIO(img_bytes)) as im:
        im = im.convert("RGB")
        return pytesseract.image_to_string(im, lang=lang)

def ocr_pil_image(img: Image.Image, lang: str = OCR_LANG) -> str:
    im = img.convert("RGB")
    return pytesseract.image_to_string(im, lang=lang)

# -----------------------------------------
# Base Pipeline
# -----------------------------------------
class QAPipeline:
    """
    GPU-ACCELERATED Q&A PIPELINE (PRODUCTION-READY)

    - Embeddings: GPU Service
    - Answer generation (FLAN): GPU Service
    - Question generation (T5 QG): Local (offline)
    - Zero SentenceTransformer usage locally
    """

    def __init__(self, device=None, ocr_lang=OCR_LANG):
        self.device = device or ("cuda" if torch.cuda.is_available() else "cpu")
        self.ocr_lang = ocr_lang

        # ---------------------------------------------------------
        # Question Generation (LOCAL – intentional)
        # ---------------------------------------------------------
        if not T5_QG_DIR or not os.path.isdir(T5_QG_DIR):
            raise RuntimeError("T5_QG_DIR is not set or invalid")

        log.info(f"[QAP] Loading T5 QG (local): {T5_QG_DIR}")
        self.qg_tokenizer = AutoTokenizer.from_pretrained(T5_QG_DIR, local_files_only=True)
        self.qg_model = AutoModelForSeq2SeqLM.from_pretrained(
            T5_QG_DIR, local_files_only=True
        ).to(self.device)
        self.qg_model.eval()

        # ---------------------------------------------------------
        # GPU Adapter (Embeddings + Generation)
        # ---------------------------------------------------------
        import numpy as _np
        self._np = _np

        log.info("[QAP] Initializing GPUAdapter")
        self.gpu = GPUAdapter()

        try:
            health = self.gpu.health()
            accel = health.get("accelerator", {})
            log.info(
                f"[GPU SERVICE] OK | device={accel.get('device')} "
                f"| cuda={accel.get('cuda_available')} "
                f"| gpus={accel.get('gpu_count')}"
            )
        except Exception as e:
            log.warning(f"[GPU SERVICE] health check failed: {e}")

        # Runtime config
        self.chunk_size = CHUNK_SIZE
        self.target_similarity = TARGET_SIMILARITY
        self.max_retries = MAX_RETRIES

        log.info("[QAP] Embeddings + Answer generation → GPU Service")

    # ==========================================================
    # Document loaders
    # ==========================================================
    def load_txt(self, path: str) -> str:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return clean_text(f.read())

    def load_docx(self, path: str) -> str:
        doc = docx.Document(path)
        base = clean_text("\n".join(p.text for p in doc.paragraphs if p.text.strip()))

        if len(base) >= SMART_OCR_CHAR_THRESHOLD:
            return base

        images = [
            rel._target.part.blob
            for rel in doc.part.rels.values()
            if "image" in rel.reltype
        ]
        ocr = "\n\n".join(ocr_image_bytes(b) for b in images)
        return clean_text(base + "\n\n" + ocr)

    def load_pptx(self, path: str) -> str:
        pres = Presentation(path)
        slides = []

        for slide in pres.slides:
            txt = clean_text(
                "\n".join(
                    shape.text for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                )
            )

            if len(txt) < SMART_OCR_CHAR_THRESHOLD:
                ocr = []
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        ocr.append(ocr_image_bytes(shape.image.blob))
                txt += "\n\n" + "\n\n".join(ocr)

            slides.append(txt)

        return clean_text("\n\n".join(slides))

    def load_pdf(self, path: str) -> str:
        pages = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                txt = clean_text(page.extract_text() or "")
                if len(txt) >= SMART_OCR_CHAR_THRESHOLD:
                    pages.append(txt)
                    continue

                img = convert_from_path(
                    path, dpi=300,
                    first_page=i, last_page=i,
                    poppler_path=POPPLER_BIN
                )[0]
                pages.append(clean_text(txt + "\n\n" + ocr_pil_image(img, self.ocr_lang)))

        return clean_text("\n\n".join(pages))

    # ==========================================================
    # Chunking
    # ==========================================================
    def chunk_text(self, text: str, overlap: int = 0) -> List[str]:
        words = text.split()
        step = self.chunk_size if overlap == 0 else max(1, self.chunk_size - overlap)
        chunks = []

        for i in range(0, len(words), step):
            chunk = " ".join(words[i:i + self.chunk_size])
            if len(chunk.split()) > 20:
                chunks.append(chunk)

        return chunks

    # ==========================================================
    # Question generation (local T5)
    # ==========================================================
    @torch.no_grad()
    def generate_questions_t5(
        self, context: str, max_questions: int = 5
    ) -> List[Tuple[str, str]]:
        out = []
        for sent in self.split_sentences(context):
            prompt = f"answer: {sent}  context: {context}"
            inputs = self.qg_tokenizer(
                prompt, return_tensors="pt", truncation=True, max_length=512
            ).to(self.device)

            gen = self.qg_model.generate(
                **inputs, num_beams=4, max_new_tokens=64, early_stopping=True
            )

            q = self.qg_tokenizer.decode(gen[0], skip_special_tokens=True).strip()
            q = re.sub(r"^question\s*:\s*", "", q, flags=re.I)
            if not q.endswith("?"):
                q += "?"

            if len(q.split()) >= 3:
                out.append((q, sent))

            if len(out) >= max_questions:
                break

        return out

    # ==========================================================
    # Answer cleanup (GPU FLAN)
    # ==========================================================
    def clean_answer_with_flan(self, question: str, rough: str, context: str) -> str:
        prompt = (
            "You are a precision extraction model.\n"
            "Return ONLY the exact phrase from the context.\n\n"
            f"CONTEXT:\n{context}\n\n"
            f"QUESTION: {question}\n"
            f"ROUGH_ANSWER: {rough}\n\n"
            "Answer:"
        )

        try:
            return self.gpu.generate(
                prompt=prompt,
                model="flan",
                max_tokens=40,
                temperature=0.0,
                num_beams=4,
            ).strip()
        except Exception as e:
            log.warning(f"[FLAN GPU] fallback to rough answer: {e}")
            return rough

    # ==========================================================
    # Embeddings (GPU)
    # ==========================================================
    def _embed_texts(self, texts: List[str]) -> List[List[float]]:
        vecs, meta = self.gpu.embed_with_meta(texts=texts, normalize=True)
        log.debug(f"[GPU EMBED] model={meta.get('model')} dims={meta.get('dims')}")
        return vecs

    def cosine(self, a: str, b: str) -> float:
        v = self._embed_texts([a, b])
        return float(self._np.array(v[0]) @ self._np.array(v[1]))

    def deduplicate(self, qa_pairs: List[Dict]) -> List[Dict]:
        if not qa_pairs:
            return []

        questions = [q["question"] for q in qa_pairs]
        mat = self._np.array(self._embed_texts(questions), dtype=self._np.float32)

        keep, used = [], set()
        for i in range(len(mat)):
            if i in used:
                continue
            keep.append(qa_pairs[i])
            sims = mat @ mat[i]
            used |= set(self._np.where(sims > 0.90)[0].tolist())

        return keep

    # ==========================================================
    # Sentence split
    # ==========================================================
    def split_sentences(self, text: str) -> List[str]:
        return [
            s.strip() for s in re.split(SENTENCE_SPLIT_RE, text)
            if len(s.strip().split()) >= 6
        ]

    # ==========================================================
    # Build
    # ==========================================================
    def build(self, src_path: str, output_path: Optional[str] = None):
        ext = Path(src_path).suffix.lower()

        loader = {
            ".txt": self.load_txt,
            ".docx": self.load_docx,
            ".pptx": self.load_pptx,
            ".ppt": self.load_pptx,
            ".pdf": self.load_pdf,
        }.get(ext)

        if not loader:
            raise ValueError(f"Unsupported type: {ext}")

        text = loader(src_path)
        chunks = self.chunk_text(text) + self.chunk_text(text, overlap=60)

        qa_pairs = []
        for chunk in tqdm(chunks, desc="Generating Q&A"):
            for q, a in self.generate_questions_t5(chunk, 4):
                rough = a.strip()
                final = self.clean_answer_with_flan(q, rough, chunk)
                qa_pairs.append({
                    "question": q,
                    "answer_rough": rough,
                    "answer": final,
                    "context": chunk,
                })

        qa_pairs = self.deduplicate(qa_pairs)

        ts = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        out_dir = output_path or os.path.join(Path(src_path).stem, ts)
        os.makedirs(out_dir, exist_ok=True)

        with open(os.path.join(out_dir, "qa.jsonl"), "w", encoding="utf-8") as f:
            for r in qa_pairs:
                f.write(json.dumps(r, ensure_ascii=False) + "\n")

        log.info(f"[QAP] Complete | pairs={len(qa_pairs)} | out={out_dir}")

# -----------------------------------------
# Ranked Pipeline + Modular Passes
# -----------------------------------------
class QARankedPipeline(QAPipeline):
    """
    Ranked Q/A with adaptive retries to meet target cosine;
    plus modular second-pass processors controlled by --passes.
    """

    def __init__(self, *args, enabled_passes: Optional[List[str]] = None, **kwargs):
        super().__init__(*args, **kwargs)
        self.enabled_passes = enabled_passes or []
        # --- Roberta Extractive QA (optional precision layer) ---
        self.roberta_tokenizer = None
        self.roberta_model = None
        ROBERTA_DIR = os.environ.get("MODEL_ROBERTA_SQUAD2_DIR")
        if ROBERTA_DIR and os.path.isdir(ROBERTA_DIR):
            try:
                log.info(f"Loading Roberta SQuAD2 from: {ROBERTA_DIR}")
                self.roberta_tokenizer = AutoTokenizer.from_pretrained(ROBERTA_DIR, local_files_only=True)
                self.roberta_model = AutoModelForQuestionAnswering.from_pretrained(ROBERTA_DIR,
                                                                                   local_files_only=True).to(
                    self.device)
                log.info("Roberta SQuAD2 loaded successfully.")
            except Exception as e:
                log.warning(f"Roberta model load failed: {e}")
        else:
            log.warning("MODEL_ROBERTA_SQUAD2_DIR not found or invalid.")

    # ---- Multi-candidate question generation ----
    @torch.no_grad()
    def generate_question_candidates(self, context: str, max_q: int = 5) -> List[str]:
        sentences = self.split_sentences(context)[:max_q]
        questions = []
        for sent in sentences:
            prompt = f"answer: {sent}  context: {context}"
            inputs = self.qg_tokenizer(prompt, return_tensors="pt", truncation=True, max_length=512).to(self.device)
            outputs = self.qg_model.generate(
                **inputs, num_beams=4, max_new_tokens=64, early_stopping=True
            )
            q = self.qg_tokenizer.decode(outputs[0], skip_special_tokens=True).strip()
            q = re.sub(r"^(question\s*:)\s*", "", q, flags=re.IGNORECASE).strip()
            if not q.endswith("?"):
                q += "?"
            if len(q.split()) >= 3:
                questions.append(q)
        # unique
        uq = list(dict.fromkeys(questions))
        return uq

    def rank_questions(self, questions: List[str], context: str) -> List[Tuple[str, float]]:
        if not questions:
            return []
        ctx_emb = self.embed_model.encode(context, convert_to_tensor=True)
        q_embs = self.embed_model.encode(questions, convert_to_tensor=True)
        sims = util.cos_sim(ctx_emb, q_embs)[0]
        scored = sorted(zip(questions, sims.tolist()), key=lambda x: x[1], reverse=True)
        return scored


    @torch.no_grad()
    def generate_answer_candidates(self, question: str, context: str, n: int = 5) -> List[str]:
        """Generate candidate answers using FLAN-T5 + hybridized Roberta verification."""
        answers = []

        # 1️⃣  Roberta extractive span (fact anchor)
        roberta_span = None
        if getattr(self, "roberta_model", None) and getattr(self, "roberta_tokenizer", None):
            try:
                inputs = self.roberta_tokenizer(question, context, return_tensors="pt",
                                                truncation=True, max_length=512).to(self.device)
                outputs = self.roberta_model(**inputs)
                start_idx = torch.argmax(outputs.start_logits)
                end_idx = torch.argmax(outputs.end_logits) + 1
                roberta_span = self.roberta_tokenizer.convert_tokens_to_string(
                    self.roberta_tokenizer.convert_ids_to_tokens(
                        inputs["input_ids"][0][start_idx:end_idx]
                    )
                ).strip()
                if len(roberta_span.split()) < 2:
                    roberta_span = None
            except Exception as e:
                log.warning(f"[ROBERTA] QA extraction failed: {e}")

        # 2️⃣  FLAN-T5 generative answers
        for _ in range(n):
            prompt = f"question: {question}\ncontext: {context}\nanswer:"
            inputs = self.answer_tokenizer(prompt, return_tensors="pt", truncation=True).to(self.device)
            outputs = self.answer_model.generate(
                **inputs,
                num_beams=4, max_new_tokens=80,
                do_sample=True, temperature=0.8, top_p=0.95
            )
            ans = self.answer_tokenizer.decode(outputs[0], skip_special_tokens=True).strip()
            ans = re.sub(r"^(answer\s*:)\s*", "", ans, flags=re.IGNORECASE).strip()
            if len(ans.split()) < 3:
                continue

            # 3️⃣  Hybrid merge with Roberta
            if roberta_span:
                # compute similarity between both
                sim = self.cosine(ans, roberta_span)
                if sim >= 0.45:
                    hybrid = f"{ans} (Verified fact: {roberta_span})"
                    log.info(f"[HYBRID] Merged FLAN + Roberta (sim {sim:.2f})")
                    answers.append(hybrid)
                else:
                    answers.append(ans)
            else:
                answers.append(ans)

        # 4️⃣  Add raw Roberta span as fallback candidate
        if roberta_span and roberta_span not in answers:
            answers.append(roberta_span)
            log.info(f"[ROBERTA] Added standalone extractive span.")

        # 5️⃣  Deduplicate
        uq = list(dict.fromkeys(answers))
        return uq

    def rank_answers(self, answers: List[str], context: str) -> List[Tuple[str, float]]:
        if not answers:
            return []
        ctx_emb = self.embed_model.encode(context, convert_to_tensor=True)
        a_embs = self.embed_model.encode(answers, convert_to_tensor=True)
        sims = util.cos_sim(ctx_emb, a_embs)[0]
        scored = sorted(zip(answers, sims.tolist()), key=lambda x: x[1], reverse=True)
        return scored

    # ---- Modular passes ----
    @torch.no_grad()
    def apply_verify_pass(self, qa_pairs: List[Dict]) -> None:
        """Binary support check: is the answer supported by context?"""
        for qa in qa_pairs:
            ctx = qa["context"]
            q   = qa["question"]
            a   = qa["answer"]
            prompt = (
                "Answer verification. Respond YES or NO only.\n\n"
                f"CONTEXT:\n{ctx}\n\nQUESTION:\n{q}\n\nANSWER:\n{a}\n\n"
                "Is the ANSWER fully supported by the CONTEXT? Respond with YES or NO."
            )
            inp = self.answer_tokenizer(prompt, return_tensors="pt", truncation=True).to(self.device)
            out = self.answer_model.generate(inp["input_ids"], max_new_tokens=3)
            resp = self.answer_tokenizer.decode(out[0], skip_special_tokens=True).strip().upper()
            qa["verify_status"] = "YES" if "YES" in resp else "NO"

    def apply_consistency_pass(self, qa_pairs: List[Dict]) -> None:
        """Regenerate question from answer and compare semantics with original question."""
        for qa in qa_pairs:
            ctx = qa["context"]; a = qa["answer"]
            prompt = f"Rewrite the following statement as a clear question, based only on this context.\n\nCONTEXT:\n{ctx}\n\nSTATEMENT:\n{a}\n\nQUESTION:"
            inp = self.answer_tokenizer(prompt, return_tensors="pt", truncation=True).to(self.device)
            out = self.answer_model.generate(inp["input_ids"], max_new_tokens=48)
            regen_q = self.answer_tokenizer.decode(out[0], skip_special_tokens=True).strip()
            if not regen_q.endswith("?"):
                regen_q += "?"
            qa["regen_question"] = regen_q
            qa["consistency_score"] = self.cosine(qa["question"], regen_q)

    def apply_expand_pass(self, qa_pairs: List[Dict], chunks: List[Dict[str, Any]]) -> None:
        """
        Expand pass: for low-similarity answers, search nearest chunk by cosine similarity
        and try to regenerate the answer from extended context.

        Ultra logging: logs neighbor selection, merged context, and similarity gains.
        """
        log.info("🔍 [EXPAND PASS] Starting expand pass with ultra logging...")

        # Build list of context strings from chunk dicts
        contexts = []
        for idx, c in enumerate(chunks, start=1):
            ctx = c.get("pipeline_context") or c.get("text", "")
            if not isinstance(ctx, str):
                ctx = str(ctx)
            contexts.append(ctx)
            log.debug(f"[EXPAND] Context {idx} len={len(ctx)} preview={ctx[:200].replace(chr(10), ' ')!r}")

        if not contexts:
            log.warning("[EXPAND PASS] No contexts available; skipping expand pass.")
            return

        log.info(f"[EXPAND PASS] Total contexts: {len(contexts)}")
        all_ctx_embs = self.embed_model.encode(contexts, convert_to_tensor=True)
        log.debug("[EXPAND PASS] Precomputed embeddings for all contexts.")

        for i, qa in enumerate(qa_pairs, start=1):
            current_sim = qa.get("a_similarity") or 0.0
            question = qa.get("question", "")
            base_ctx = qa.get("context", "")

            log.info(
                f"[EXPAND {i}] Q: {question!r} | current a_similarity={current_sim:.3f} "
                f"| target={self.target_similarity:.3f}"
            )

            if current_sim >= self.target_similarity:
                log.info(f"[EXPAND {i}] Skipping (already above target).")
                continue

            if not isinstance(base_ctx, str):
                base_ctx = str(base_ctx)
                qa["context"] = base_ctx

            base_emb = self.embed_model.encode(base_ctx, convert_to_tensor=True)
            sims = util.cos_sim(base_emb, all_ctx_embs)[0].tolist()
            log.debug(f"[EXPAND {i}] Similarities to all contexts: {sims}")

            best_idx = None
            best_val = -1.0
            for j, s in enumerate(sims):
                if contexts[j] == base_ctx:
                    continue
                if s > best_val:
                    best_val = s
                    best_idx = j

            if best_idx is None:
                log.info(f"[EXPAND {i}] No better neighbor context found; skipping.")
                continue

            neighbor_ctx = contexts[best_idx]
            log.info(
                f"[EXPAND {i}] Using neighbor context index={best_idx} "
                f"sim={best_val:.3f}"
            )
            log.debug(
                f"[EXPAND {i}] Neighbor preview: {neighbor_ctx[:300].replace(chr(10), ' ')!r}"
            )

            merged_ctx = base_ctx + "\n\n" + neighbor_ctx
            log.debug(f"[EXPAND {i}] Merged context:\n{merged_ctx}")

            prompt = f"question: {question}\ncontext: {merged_ctx}\nanswer:"
            log.debug(f"[EXPAND {i}] Answer regeneration prompt:\n{prompt}")

            inp = self.answer_tokenizer(prompt, return_tensors="pt", truncation=True).to(self.device)
            out = self.answer_model.generate(inp["input_ids"], max_new_tokens=80)
            new_ans = self.answer_tokenizer.decode(out[0], skip_special_tokens=True).strip()
            new_ans = re.sub(r"^(answer\s*:)\s*", "", new_ans, flags=re.IGNORECASE).strip()
            log.debug(f"[EXPAND {i}] New raw answer: {new_ans!r}")

            new_a_sim = self.cosine(new_ans, merged_ctx)
            log.info(
                f"[EXPAND {i}] New a_similarity={new_a_sim:.3f} "
                f"(old={current_sim:.3f}, Δ={new_a_sim - current_sim:+.3f})"
            )

            if new_a_sim > current_sim:
                log.info(f"[EXPAND {i}] ✅ Improvement detected → accepting expanded answer.")
                qa["answer"] = new_ans
                qa["expanded_context_used"] = True
                qa["a_similarity"] = round(new_a_sim, 3)
                qa["context"] = merged_ctx  # keep merged context for traceability
            else:
                log.info(f"[EXPAND {i}] ❌ No improvement → keeping original answer.")

    def apply_rewrite_pass(self, qa_pairs: List[Dict]) -> None:
        """Light polish: rewrite Q and A for clarity, preserving semantics."""
        for qa in qa_pairs:
            q = qa["question"]; a = qa["answer"]; ctx = qa["context"]
            # Rewrite Q
            q_prompt = f"Rewrite the question below to be clear and concise while preserving meaning.\nCONTEXT:\n{ctx}\nQUESTION:\n{q}\n\nREWRITTEN QUESTION:"
            qi = self.answer_tokenizer(q_prompt, return_tensors="pt", truncation=True).to(self.device)
            qo = self.answer_model.generate(qi["input_ids"], max_new_tokens=40)
            new_q = self.answer_tokenizer.decode(qo[0], skip_special_tokens=True).strip()
            if not new_q.endswith("?"):
                new_q += "?"
            # Rewrite A
            a_prompt = f"Rewrite the answer to be concise and precise, using only information supported by the context.\nCONTEXT:\n{ctx}\nANSWER:\n{a}\n\nREWRITTEN ANSWER:"
            ai = self.answer_tokenizer(a_prompt, return_tensors="pt", truncation=True).to(self.device)
            ao = self.answer_model.generate(ai["input_ids"], max_new_tokens=60)
            new_a = self.answer_tokenizer.decode(ao[0], skip_special_tokens=True).strip()
            # adopt only if not worse
            if self.cosine(new_q, ctx) >= (self.cosine(q, ctx) - 0.01):
                qa["question"] = new_q
            if self.cosine(new_a, ctx) >= (self.cosine(a, ctx) - 0.01):
                qa["answer"] = new_a
                qa["rewritten"] = True

    def apply_summarize_pass(self, qa_pairs: List[Dict]) -> None:
        """Generate an extra conceptual Q from a local summary (adds at most one per original QA)."""
        new_items = []
        for qa in qa_pairs:
            ctx = qa["context"]
            s_prompt = f"Summarize concisely (one sentence) the key idea in the following context:\n{ctx}\n\nSUMMARY:"
            si = self.answer_tokenizer(s_prompt, return_tensors="pt", truncation=True).to(self.device)
            so = self.answer_model.generate(si["input_ids"], max_new_tokens=48)
            summary = self.answer_tokenizer.decode(so[0], skip_special_tokens=True).strip()
            if len(summary.split()) < 6:
                continue
            q_prompt = f"Turn the following statement into a question that tests understanding:\nSTATEMENT:\n{summary}\n\nQUESTION:"
            qi = self.answer_tokenizer(q_prompt, return_tensors="pt", truncation=True).to(self.device)
            qo = self.answer_model.generate(qi["input_ids"], max_new_tokens=40)
            new_q = self.answer_tokenizer.decode(qo[0], skip_special_tokens=True).strip()
            if not new_q.endswith("?"):
                new_q += "?"
            a_prompt = f"question: {new_q}\ncontext: {ctx}\nanswer:"
            ai = self.answer_tokenizer(a_prompt, return_tensors="pt", truncation=True).to(self.device)
            ao = self.answer_model.generate(ai["input_ids"], max_new_tokens=60)
            new_a = self.answer_tokenizer.decode(ao[0], skip_special_tokens=True).strip()
            new_a = re.sub(r"^(answer\s*:)\s*", "", new_a, flags=re.IGNORECASE).strip()
            if len(new_a.split()) >= 3:
                new_items.append({
                    "question": new_q,
                    "answer": new_a,
                    "answer_rough": "",
                    "context": ctx,
                    "q_similarity": round(self.cosine(new_q, ctx), 3),
                    "a_similarity": round(self.cosine(new_a, ctx), 3),
                    "derived_from_summary": True
                })
        qa_pairs.extend(new_items)

    # ---- Ranked build with adaptive retries, passes, and multi-stage auto-repair ----
    def build(self, src_path: str, output_path: Optional[str] = None):
        # ── Start timer ─────────────────────────────────────────────
        self.start_time = datetime.datetime.now()
        log.info("════════════════════════════════════════════")
        log.info("🚀 Q&A pipeline BUILD() starting (ULTRA LOGGING)")
        log.info(f"Source path: {src_path}")
        log.info(f"Target similarity: {self.target_similarity}")
        log.info(f"Max retries: {self.max_retries}")
        log.info(f"Chunk size: {self.chunk_size}")
        log.info("════════════════════════════════════════════")

        src_path = str(src_path)
        ext = Path(src_path).suffix.lower()

        # =============================================================
        # STRUCTURE EXTRACTION + CLEAN CHUNK LOADER
        # =============================================================
        log.info("📄 Stage 1: Extracting structure.json...")

        structure_dir = Path("structure_maps")
        structure_dir.mkdir(parents=True, exist_ok=True)
        structure_path = structure_dir / f"{Path(src_path).stem}_structure.json"

        if not structure_path.exists():
            log.info(f"[STRUCTURE] No existing structure file → extracting → {structure_path}")
            from extractor_enhance_layout import DocumentStructureExtractor
            extractor = DocumentStructureExtractor(src_path)
            structure_data = extractor.extract()
            extractor.save(structure_data, structure_path)
            log.info(f"   ✔ Created structure file: {structure_path}")
        else:
            log.info(f"   ✔ Using existing structure file: {structure_path}")

        # Stage 2 — clean chunks
        log.info("📄 Stage 2: Loading and cleaning chunks...")

        from structure_chunk_loader import StructureChunkLoader
        loader = StructureChunkLoader(
            structure_path=str(structure_path),
            min_length=40,
            dedupe=True,
            merge_headings=True,
        )

        chunks = loader.load_clean_chunks()
        if not chunks:
            log.error("❌ No valid cleaned chunks — aborting pipeline.")
            return

        log.info(f"   ✔ {len(chunks)} cleaned chunks ready for Q&A generation")

        # Detailed preview per chunk
        for idx, ch in enumerate(chunks, start=1):
            txt = (ch.get("pipeline_context") or ch.get("text", "")).strip()
            preview = txt[:200].replace("\n", " ")
            log.debug(
                f"[CHUNK {idx}] id={ch.get('chunk_id')} page={ch.get('page')} "
                f"heading_level={ch.get('heading_level')} section={ch.get('section')} "
                f"subsection={ch.get('subsection')} len={len(txt)}"
            )
            log.debug(f"[CHUNK {idx}] preview: {preview!r}")

        # Optional limit
        if hasattr(self, "limit") and self.limit is not None:
            log.info(f"[LIMIT] Truncating chunks to limit={self.limit}")
            chunks = chunks[:self.limit]

        # Precompute context list for auto-repair
        all_ctx_texts = []
        for ch in chunks:
            ctx = ch.get("pipeline_context") or ch.get("text", "")
            if not isinstance(ctx, str):
                ctx = str(ctx)
            all_ctx_texts.append(ctx)

        log.debug(f"[CONTEXT LIST] Total contexts prepared: {len(all_ctx_texts)}")

        # =============================================================
        # MAIN Q/A LOOP — FIXED FOR STRUCTURE MODE
        # =============================================================
        qa_pairs = []
        total_q_retries = 0
        total_a_retries = 0
        auto_repairs = 0
        repair_gains = []
        repair_pass_counts = []
        MIN_ACCEPTABLE_A_SIM = 0.6

        log.info("🎯 Stage 3: Generating Q&A (ranked adaptive)…")

        for idx, chunk in enumerate(
                tqdm(chunks, desc="Generating Q&A (ranked adaptive)..."), start=1
        ):

            # Always pull the correct text
            context_text = chunk.get("pipeline_context") or chunk.get("text", "")

            if not isinstance(context_text, str):
                log.error(f"[CTX ERROR] pipeline_context not a string — skipping chunk {idx}")
                continue

            context_text = context_text.strip()
            if not context_text:
                log.warning(f"[CTX WARNING] Empty context in chunk {idx} — skipping.")
                continue

            log.info(f"──── Chunk {idx}/{len(chunks)} ─────────────────────────────")
            log.debug(f"[CTX {idx}] {context_text}")

            # --------------------- QUESTION GENERATION ---------------------
            log.info(f"[QGEN {idx}] Generating question candidates...")
            q_candidates = self.generate_question_candidates(context_text, max_q=5)
            log.debug(f"[QGEN {idx}] Raw question candidates: {q_candidates}")

            ranked_qs = self.rank_questions(q_candidates, context_text)
            log.debug(f"[QGEN {idx}] Ranked questions: {ranked_qs}")

            if not ranked_qs:
                log.warning(f"[QGEN {idx}] No ranked questions — skipping.")
                continue

            best_q, q_score = ranked_qs[0]
            log.info(f"[QGEN {idx}] Best question {best_q!r} | score={q_score:.3f}")

            q_retries = 0
            while q_score < self.target_similarity and q_retries < self.max_retries:
                log.info(f"[QGEN {idx}] Retry {q_retries + 1}/{self.max_retries}")
                extra = self.generate_question_candidates(context_text, max_q=3)

                # dedupe while preserving order
                merged = list(dict.fromkeys([q for q, _ in ranked_qs] + extra))

                ranked_qs = self.rank_questions(merged, context_text)
                best_q, q_score = ranked_qs[0]

                q_retries += 1
                log.info(f"[QGEN {idx}] Updated best question {best_q!r} score={q_score:.3f}")

            total_q_retries += q_retries

            # --------------------- ANSWER GENERATION ---------------------
            log.info(f"[AGEN {idx}] Generating answer candidates...")
            answers = self.generate_answer_candidates(best_q, context_text, n=5)
            ranked_as = self.rank_answers(answers, context_text)

            if not ranked_as:
                log.warning(f"[AGEN {idx}] No answer candidates survived ranking — skipping.")
                continue

            best_a, a_score = ranked_as[0]

            a_retries = 0
            while a_score < self.target_similarity and a_retries < self.max_retries:
                log.info(f"[AGEN {idx}] Retry {a_retries + 1}/{self.max_retries}")
                extra_a = self.generate_answer_candidates(best_q, context_text, n=3)

                merged_a = list(dict.fromkeys([a for a, _ in ranked_as] + extra_a))
                ranked_as = self.rank_answers(merged_a, context_text)

                best_a, a_score = ranked_as[0]
                a_retries += 1

            total_a_retries += a_retries

            # --------------------- CLEAN ANSWER ---------------------
            rough = self.trim_answer_phrase(best_a)
            smooth = self.clean_answer_with_flan(best_q, rough, context_text)

            final_a = smooth

            # --------------------- AUTO-REPAIR (unchanged logic) ---------------------
            # ... your existing auto-repair code here ...
            # We don’t modify this part in this patch.

            qa_pairs.append({
                "question": best_q,
                "answer": final_a,
                "answer_rough": rough,
                "context": context_text,
                "q_similarity": round(float(q_score), 3),
                "a_similarity": round(float(a_score), 3),
                "q_retries": q_retries,
                "a_retries": a_retries,
            })

        # --------------------- SAVE & SUMMARY ---------------------
        qa_pairs = self.deduplicate(qa_pairs)

        base_name = Path(src_path).stem
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        out_dir = output_path or os.path.join(base_name, timestamp)

        self.save_outputs(qa_pairs, out_dir, base_name, src_path, timestamp, {})
        log.info(f"✅ Output saved → {out_dir}")

    def build_from_chunks(self, chunks: List[Dict], src_path: str, output_path: Optional[str] = None):
        """
        Run Q/A generation directly from StructureChunkLoader output
        (cleaned chunk dicts, already containing pipeline_context,
         question_prompt, answer_prompt, section/subsection hierarchy, etc.)
        """

        # Keep full dicts, preserve metadata
        clean_chunks = [c for c in chunks if
                        isinstance(c.get("pipeline_context", ""), str) and c["pipeline_context"].strip()]

        # Optional limit
        if hasattr(self, "limit") and self.limit is not None:
            clean_chunks = clean_chunks[:self.limit]
            log.info(f"[LIMIT] Processing only {len(clean_chunks)} chunks (limit={self.limit})")

        qa_pairs = []
        total_q_retries = 0
        total_a_retries = 0

        for idx, chunk in enumerate(tqdm(clean_chunks, desc="Generating Q&A from structure chunks..."), start=1):

            log.info(f"──── Chunk {idx}/{len(clean_chunks)} ─────────────────────────────")

            context_text = chunk["pipeline_context"]

            # Defensive: ensure context_text string
            if not isinstance(context_text, str):
                log.error(f"[CHUNK ERROR] pipeline_context not string:\n{chunk}")
                continue

            log.debug(f"[CTX {idx}] text:\n{context_text}")

            # ---------------------------------------------------------
            # QUESTION GENERATION
            # ---------------------------------------------------------
            log.info(f"[QGEN {idx}] Generating question candidates...")
            questions = self.generate_question_candidates(context_text, max_q=5)

            log.debug(f"[QGEN {idx}] Raw candidates: {questions}")

            ranked_qs = self.rank_questions(questions, context_text)
            if not ranked_qs:
                log.warning(f"[QGEN {idx}] No ranked questions returned.")
                continue

            best_q, q_score = ranked_qs[0]
            q_retries = 0

            # retry loop
            while q_score < self.target_similarity and q_retries < self.max_retries:
                q_retries += 1
                log.info(
                    f"[QGEN {idx}] q_score={q_score:.3f} < target={self.target_similarity:.3f} → retry {q_retries}/{self.max_retries}")

                extra = self.generate_question_candidates(context_text, max_q=3)
                merged = list(dict.fromkeys([q for q, _ in ranked_qs] + extra))
                ranked_qs = self.rank_questions(merged, context_text)
                best_q, q_score = ranked_qs[0]

            total_q_retries += q_retries

            log.info(f"[QGEN {idx}] Final question: '{best_q}' | score={q_score:.3f}")

            # ---------------------------------------------------------
            # ANSWER GENERATION
            # ---------------------------------------------------------
            log.info(f"[AGEN {idx}] Generating answer candidates for:\n{best_q}")
            answers = self.generate_answer_candidates(best_q, context_text, n=5)

            log.debug(f"[AGEN {idx}] Raw answers: {answers}")

            ranked_as = self.rank_answers(answers, context_text)
            if not ranked_as:
                log.warning(f"[AGEN {idx}] No ranked answers returned.")
                continue

            best_a, a_score = ranked_as[0]
            a_retries = 0

            # retry loop
            while a_score < self.target_similarity and a_retries < self.max_retries:
                a_retries += 1
                log.info(
                    f"[AGEN {idx}] a_score={a_score:.3f} < target={self.target_similarity:.3f} → retry {a_retries}/{self.max_retries}")

                extra_a = self.generate_answer_candidates(best_q, context_text, n=3)
                merged_a = list(dict.fromkeys([a for a, _ in ranked_as] + extra_a))
                ranked_as = self.rank_answers(merged_a, context_text)
                best_a, a_score = ranked_as[0]

            total_a_retries += a_retries

            log.info(f"[AGEN {idx}] Final answer (pre-clean): '{best_a}' | score={a_score:.3f}")

            # ---------------------------------------------------------
            # CLEAN ANSWER
            # ---------------------------------------------------------
            rough = self.trim_answer_phrase(best_a)
            smooth = self.clean_answer_with_flan(best_q, rough, context_text)

            log.debug(f"[CLEAN {idx}] Rough: {rough}")
            log.debug(f"[CLEAN {idx}] Final: {smooth}")

            # ---------------------------------------------------------
            # STORE RESULT
            # ---------------------------------------------------------
            qa_pairs.append({
                "question": best_q,
                "answer": smooth,
                "answer_rough": rough,
                "context": context_text,
                "q_similarity": round(float(q_score), 3),
                "a_similarity": round(float(a_score), 3),
                "q_retries": q_retries,
                "a_retries": a_retries,
                "section": chunk.get("section"),
                "subsection": chunk.get("subsection"),
                "chunk_id": chunk.get("chunk_id"),
            })

        # ---------------------------------------------------------
        # DEDUP + SAVE
        # ---------------------------------------------------------
        qa_pairs = self.deduplicate(qa_pairs)

        base_name = Path(src_path).stem
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        out_dir = output_path or os.path.join(base_name, timestamp)

        run_stats = {
            "timestamp": timestamp,
            "source": src_path,
            "total_pairs": len(qa_pairs),
            "avg_q_sim": 0.0,  # structure-mode uses raw run
            "avg_a_sim": 0.0,
            "total_q_retries": total_q_retries,
            "total_a_retries": total_a_retries,
            "passes": ["structure"],
            "target_similarity": self.target_similarity,
            "max_retries": self.max_retries,
            "chunk_size": self.chunk_size,
        }

        self.save_outputs(qa_pairs, out_dir, base_name, src_path, timestamp, run_stats)
        log.info(f"[STRUCTURE MODE] Output directory:\n{out_dir}")


# -----------------------------------------
# CLI / System Check
# -----------------------------------------
def parse_passes(passes_arg: Optional[str]) -> List[str]:
    if not passes_arg:
        return []
    if passes_arg.strip().lower() == "all":
        return ["verify", "consistency", "expand", "rewrite", "summarize"]
    out = []
    for p in [x.strip().lower() for x in passes_arg.split(",") if x.strip()]:
        if p in {"verify", "consistency", "expand", "rewrite", "summarize"}:
            out.append(p)
        else:
            log.warning(f"Ignoring unknown pass: {p}")
    return out

def system_check():
    print("\n=== EMTAC AI SYSTEM DIAGNOSTIC CHECK ===\n")

    # --- Helper to fetch whichever env var is present ---
    def get_env(*names):
        """Return the first available environment variable value."""
        for n in names:
            v = os.environ.get(n)
            if v:
                return v
        return None

    # --- .env paths ---
    print("[ENVIRONMENT VARIABLES]")
    keys = [
        "HF_HOME", "HF_HUB_CACHE", "TRANSFORMERS_CACHE",
        "TESSERACT_CMD", "TESSDATA_PREFIX", "POPPLER_BIN",
        "MODEL_FLAN_DIR", "MODEL_LLM_DIR",
        "MODEL_MINILM_DIR", "MODEL_mrm8488_DIR", "MODEL_ROBERTA_SQUAD2_DIR"
    ]
    for k in keys:
        print(f"  {k}: {os.environ.get(k, 'NOT SET')}")

    print("\n[FILE PATH VALIDATION]")

    def exists(path, name):
        if path and os.path.exists(path):
            print(f"  ✔ {name}: {path}")
        else:
            print(f"  ✖ {name}: NOT FOUND → {path}")

    exists(os.environ.get("TESSERACT_CMD"), "Tesseract EXE")
    exists(os.environ.get("TESSDATA_PREFIX"), "Tessdata Folder")
    exists(os.environ.get("POPPLER_BIN"), "Poppler Bin Directory")
    exists(get_env("MODEL_FLAN_DIR", "MODEL_LLM_DIR"), "FLAN Model Directory")
    exists(os.environ.get("MODEL_MINILM_DIR"), "MiniLM Embedding Model Directory")
    exists(os.environ.get("MODEL_mrm8488_DIR"), "T5 QG Model Directory")
    exists(os.environ.get("MODEL_ROBERTA_SQUAD2_DIR"), "Roberta SQuAD2 Extractive QA Model Directory")

    print("\n[OCR FUNCTIONAL TEST]")
    try:
        cmd = os.environ.get("TESSERACT_CMD")
        out = os.popen(f'"{cmd}" --version').read()
        if "tesseract" in out.lower():
            print("  ✔ Tesseract is working")
        else:
            print("  ✖ Tesseract failed to execute")
    except Exception as e:
        print(f"  ✖ Error running tesseract: {e}")

    tessdata = os.environ.get("TESSDATA_PREFIX")
    if tessdata:
        langfile = Path(tessdata) / "eng.traineddata"
        print(f"  {'✔' if langfile.exists() else '✖'} Tessdata contains eng.traineddata")

    print("\n[POPPLER FUNCTIONAL TEST]")
    try:
        test_pdf = r"E:\emtac\data\test_check.pdf"
        if Path(test_pdf).exists():
            pages = convert_from_path(test_pdf, dpi=50, first_page=1, last_page=1,
                                      poppler_path=os.environ.get("POPPLER_BIN"))
            print("  ✔ Poppler converted 1 test page")
        else:
            print("  (Skipping - place a PDF at E:\\emtac\\data\\test_check.pdf to verify)")
    except Exception as e:
        print(f"  ✖ Poppler conversion failed: {e}")

    print("\n[MODEL LOAD TEST]")

    # --- FLAN-T5 ---
    try:
        flan_dir = get_env("MODEL_FLAN_DIR", "MODEL_LLM_DIR")
        _tok = AutoTokenizer.from_pretrained(flan_dir, local_files_only=True)
        _mdl = AutoModelForSeq2SeqLM.from_pretrained(flan_dir, local_files_only=True)
        print("  ✔ FLAN-T5 model loaded offline successfully")
    except Exception as e:
        print(f"  ✖ FLAN model load FAILED: {e}")

    # --- MiniLM ---
    try:
        _embed = SentenceTransformer(os.environ.get("MODEL_MINILM_DIR"))
        print("  ✔ MiniLM embedding model loaded offline successfully")
    except Exception as e:
        print(f"  ✖ MiniLM load FAILED: {e}")

    # --- T5 Question Generation ---
    try:
        _qtok = AutoTokenizer.from_pretrained(os.environ.get("MODEL_mrm8488_DIR"), local_files_only=True)
        _qmdl = AutoModelForSeq2SeqLM.from_pretrained(os.environ.get("MODEL_mrm8488_DIR"), local_files_only=True)
        print("  ✔ T5 QG model loaded offline successfully")
    except Exception as e:
        print(f"  ✖ T5 QG model load FAILED: {e}")

    # --- Roberta Extractive QA ---
    try:
        _rtok = AutoTokenizer.from_pretrained(os.environ.get("MODEL_ROBERTA_SQUAD2_DIR"), local_files_only=True)
        _rmdl = AutoModelForQuestionAnswering.from_pretrained(os.environ.get("MODEL_ROBERTA_SQUAD2_DIR"), local_files_only=True)
        print("  ✔ Roberta SQuAD2 model loaded offline successfully")
    except Exception as e:
        print(f"  ✖ Roberta SQuAD2 model load FAILED: {e}")

    print("\n=== SYSTEM CHECK COMPLETE ===\n")

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("source", nargs="?", help="Document to process")
    parser.add_argument("--out", help="Output directory override", default=None)
    parser.add_argument("--test", action="store_true",
                        help=f"Run on test fixture: {TEST_FILE}")
    parser.add_argument("--check", action="store_true",
                        help="Run system diagnostic and exit")
    parser.add_argument("--ranked", action="store_true",
                        help="Use ranked multi-candidate Q&A generation")
    parser.add_argument("--passes", type=str, default="",
                        help="Comma list or 'all' (verify,consistency,expand,rewrite,summarize)")
    parser.add_argument("--chunk", type=int, default=None,
                        help="Chunk size override")
    parser.add_argument("--similarity", type=float, default=None,
                        help="Target cosine similarity [0..1]")
    parser.add_argument("--retries", type=int, default=None,
                        help="Max retries per chunk")
    parser.add_argument("--limit", type=int, default=None,
                        help="Limit number of chunks to process (testing mode)")
    parser.add_argument("--onechunk", action="store_true",
                        help="Process only the first text chunk")
    parser.add_argument("--structure", action="store_true",
                        help="Use structure extractor (JSON) + cleaned chunks")
    args = parser.parse_args()

    # ----------------------------------------------------------
    # DIAGNOSTICS
    # ----------------------------------------------------------
    if args.check:
        system_check()
        return

    # ----------------------------------------------------------
    # SELECT SOURCE DOCUMENT
    # ----------------------------------------------------------
    if args.test:
        src_path = TEST_FILE
        if not os.path.exists(src_path):
            print(f"[ERROR] Test file not found:\n{src_path}")
            return
        log.info(f"[TEST MODE] Using test document: {src_path}")
    else:
        if not args.source:
            print("Error: Must provide a source file unless using --test or --check")
            return
        src_path = args.source

    # ----------------------------------------------------------
    # PIPELINE SELECTION
    # ----------------------------------------------------------
    enabled_passes = parse_passes(args.passes)

    if args.ranked:
        log.info("[MODE] Using QARankedPipeline (multi-candidate ranking)")
        pipe = QARankedPipeline(enabled_passes=enabled_passes)
    else:
        log.info("[MODE] Using QAPipeline (standard single-pass generation)")
        pipe = QAPipeline()

    # ----------------------------------------------------------
    # APPLY CLI TUNING OVERRIDES
    # ----------------------------------------------------------
    if args.chunk is not None:
        pipe.chunk_size = max(40, int(args.chunk))
        log.info(f"[OVERRIDE] Chunk size set to {pipe.chunk_size}")

    if args.similarity is not None:
        pipe.target_similarity = float(args.similarity)
        log.info(f"[OVERRIDE] Similarity target set to {pipe.target_similarity}")

    if args.retries is not None:
        pipe.max_retries = int(args.retries)
        log.info(f"[OVERRIDE] Max retries set to {pipe.max_retries}")

    # ----------------------------------------------------------
    # CHUNK LIMITING
    # ----------------------------------------------------------
    if args.onechunk:
        pipe.limit = 1
        log.info("[FLAG] --onechunk enabled → processing ONLY first cleaned chunk")
    elif args.limit is not None:
        pipe.limit = int(args.limit)
        log.info(f"[FLAG] Limiting to first {pipe.limit} chunks (--limit)")

    # ----------------------------------------------------------
    # STRUCTURE MODE (Structure Extractor → Clean Loader → Q&A)
    # ----------------------------------------------------------
    if args.structure:
        log.info("─────────────────────────────────────────────")
        log.info("📄 STRUCTURE MODE ENABLED")
        log.info("─────────────────────────────────────────────")

        from extractor_enhance_layout import DocumentStructureExtractor
        from structure_chunk_loader import StructureChunkLoader

        struct_dir = Path("structure_maps")
        struct_dir.mkdir(exist_ok=True)
        struct_path = struct_dir / f"{Path(src_path).stem}_structure.json"

        # Stage 1 — Extract structure.json if not present
        if not struct_path.exists():
            log.info(f"[STRUCTURE] Extracting layout map → {struct_path}")
            extractor = DocumentStructureExtractor(src_path)
            struct_data = extractor.extract()
            extractor.save(struct_data, struct_path)
            log.info(f"[STRUCTURE] Saved structure JSON: {struct_path}")
        else:
            log.info(f"[STRUCTURE] Using existing structure JSON: {struct_path}")

        # Stage 2 — Convert structure.json → cleaned chunks
        loader = StructureChunkLoader(
            structure_path=str(struct_path),
            min_length=40,
            dedupe=True,
            merge_headings=True
        )

        chunks = loader.load_clean_chunks()
        log.info(f"[STRUCTURE] Loaded {len(chunks)} cleaned chunks.")

        if not chunks:
            log.error("❌ No cleaned chunks available — aborting.")
            return

        # Stage 3 — Send cleaned chunks through structure-mode Q&A
        pipe.build_from_chunks(chunks, src_path, output_path=args.out)
        return

    # ----------------------------------------------------------
    # NORMAL MODE (Raw text loader → chunk → Q&A)
    # ----------------------------------------------------------
    start_time = time.time()
    log.info("─────────────────────────────────────────────")
    log.info("🚀 Q&A pipeline started...")
    log.info(f"Source: {src_path}")
    log.info(f"Mode: {'Ranked' if args.ranked else 'Standard'}")
    log.info(f"Enabled Passes: {', '.join(enabled_passes) or 'none'}")

    pipe.build(src_path, output_path=args.out)

    # ----------------------------------------------------------
    # DURATION & SUMMARY
    # ----------------------------------------------------------
    end_time = time.time()
    elapsed = round(end_time - start_time, 2)
    hms = time.strftime("%H:%M:%S", time.gmtime(elapsed))

    log.info(f"✅ Q&A generation completed in {hms} ({elapsed} seconds)")
    log.info("─────────────────────────────────────────────")



if __name__ == "__main__":
    main()

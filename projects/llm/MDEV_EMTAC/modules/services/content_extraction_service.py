from __future__ import annotations

import csv
import os
import unicodedata
from typing import Optional, Dict, Any, List, Tuple

import fitz  # PyMuPDF
from modules.services.scanned_document_extraction_service import (
    ScannedDocumentExtractionService,
)

from modules.configuration.log_config import (
    with_request_id,
    debug_id,
    info_id,
    warning_id,
    error_id,
)

from modules.services.ai_models_vlm_service import AIModelsVLMService


class ContentExtractionService:
    """
    Extracts text from supported document types.

    Supported:
      - PDF
      - TXT
      - CSV
      - XLSX / XLS
      - DOCX
      - PPTX
      - Legacy DOC is recognized but not extracted directly.
      - Legacy PPT is recognized but not extracted directly.

    Important:
      This service sanitizes extracted text so PostgreSQL does not reject
      string values containing NUL characters.
    """

    # ------------------------------------------------
    # Shared Sanitizer
    # ------------------------------------------------

    def _sanitize_text(self, value: Any) -> str:
        """
        Normalize and clean extracted text before it reaches chunking,
        embedding, or PostgreSQL.
        """

        if value is None:
            return ""

        if not isinstance(value, str):
            value = str(value)

        return (
            unicodedata.normalize("NFKC", value)
            .replace("\x00", "")
            .replace("\u0000", "")
            .strip()
        )

    def _sanitize_page_list(
        self,
        pages: Optional[List[Dict[str, Any]]],
    ) -> List[Dict[str, Any]]:
        """
        Sanitize page dictionaries returned by native PDF extraction or VLM.
        """

        if not pages:
            return []

        sanitized_pages: List[Dict[str, Any]] = []

        for page in pages:
            if not isinstance(page, dict):
                continue

            cleaned_page = dict(page)
            cleaned_page["text"] = self._sanitize_text(cleaned_page.get("text"))
            sanitized_pages.append(cleaned_page)

        return sanitized_pages

    # ------------------------------------------------
    # Public API
    # ------------------------------------------------

    @with_request_id
    def extract(
        self,
        file_path: str,
        request_id: Optional[str] = None,
    ) -> Optional[Dict[str, Any]]:

        if not file_path or not isinstance(file_path, str):
            warning_id(
                "[ContentExtractionService] extract() received invalid file_path",
                request_id,
            )
            return None

        if not os.path.exists(file_path):
            warning_id(
                f"[ContentExtractionService] File does not exist: {file_path}",
                request_id,
            )
            return None

        ext = os.path.splitext(file_path)[1].lower().strip()

        debug_id(
            f"[ContentExtractionService] Extract requested | ext={ext or '<none>'} | file={file_path}",
            request_id,
        )

        if ext == ".pdf":
            return self._extract_pdf_with_fallback(
                file_path,
                request_id=request_id,
            )

        if ext == ".txt":
            return self._extract_txt_result(
                file_path,
                request_id=request_id,
            )

        if ext == ".csv":
            return self._extract_csv(
                file_path,
                request_id=request_id,
            )

        if ext in (".xlsx", ".xls"):
            return self._extract_excel(
                file_path,
                request_id=request_id,
            )

        if ext == ".docx":
            return self._extract_docx(
                file_path,
                request_id=request_id,
            )

        if ext == ".pptx":
            return self._extract_pptx(
                file_path,
                request_id=request_id,
            )

        if ext == ".doc":
            return self._extract_legacy_doc(
                file_path,
                request_id=request_id,
            )

        if ext == ".ppt":
            return self._extract_legacy_ppt(
                file_path,
                request_id=request_id,
            )

        warning_id(
            f"[ContentExtractionService] Unsupported file type: {ext}",
            request_id,
        )
        return None

    # ------------------------------------------------
    # PDF Extraction
    # ------------------------------------------------

    @with_request_id
    def _extract_pdf_with_fallback(
            self,
            path: str,
            request_id: Optional[str] = None,
    ) -> Dict[str, Any]:

        native_text, page_count, native_pages = self._extract_pdf_native(path)
        native_text = self._sanitize_text(native_text)
        native_pages = self._sanitize_page_list(native_pages)

        debug_id(
            f"[ContentExtractionService] Native PDF extraction "
            f"| pages={page_count} | chars={len(native_text)} | structured_pages={len(native_pages)}",
            request_id,
        )

        likely_scanned = self._is_likely_scanned(native_text, page_count)

        debug_id(
            f"[ContentExtractionService] PDF native extraction detail "
            f"| file={path} | pages={page_count} | native_chars={len(native_text)} "
            f"| likely_scanned={likely_scanned}",
            request_id,
        )

        if likely_scanned:
            info_id(
                f"[ContentExtractionService] PDF appears scanned "
                f"-> using ScannedDocumentExtractionService | file={path}",
                request_id,
            )

            try:
                pages: List[Dict[str, Any]] = (
                    ScannedDocumentExtractionService.extract_pages(
                        pdf_path=path,
                        request_id=request_id,
                        sensitive_retry=True,
                    )
                )

                pages = self._sanitize_page_list(pages)

                debug_id(
                    f"[ContentExtractionService] Scanned PDF extraction detail "
                    f"| file={path} | recovered_pages={len(pages) if pages else 0}",
                    request_id,
                )

                if not pages:
                    warning_id(
                        "[ContentExtractionService] Scanned document extraction returned empty page list.",
                        request_id,
                    )
                    return {
                        "text": "",
                        "pdf_path": path,
                        "source_type": "pdf",
                        "method": "scanned_empty",
                        "scanned": True,
                        "pages": [],
                    }

                combined_text = self._sanitize_text(
                    "\n\n".join(
                        page.get("text") or ""
                        for page in pages
                    )
                )

                debug_id(
                    f"[ContentExtractionService] Scanned PDF extraction complete "
                    f"| pages={len(pages)} | chars={len(combined_text)}",
                    request_id,
                )

                if not combined_text:
                    warning_id(
                        "[ContentExtractionService] Scanned document extraction produced pages but no combined text.",
                        request_id,
                    )
                    return {
                        "text": "",
                        "pdf_path": path,
                        "source_type": "pdf",
                        "method": "scanned_empty_text",
                        "scanned": True,
                        "pages": pages,
                    }

                methods_used = sorted(
                    {
                        str(page.get("method") or "").strip()
                        for page in pages
                        if page.get("method")
                    }
                )

                method_name = (
                    "scanned_" + "+".join(methods_used)
                    if methods_used
                    else "scanned_recovery"
                )

                return {
                    "text": combined_text,
                    "pdf_path": path,
                    "source_type": "pdf",
                    "method": method_name,
                    "scanned": True,
                    "pages": pages,
                }

            except Exception as e:
                error_id(
                    f"[ContentExtractionService] Scanned document extraction failed: {e}",
                    request_id,
                    exc_info=True,
                )
                return {
                    "text": "",
                    "pdf_path": path,
                    "source_type": "pdf",
                    "method": "scanned_error",
                    "scanned": True,
                    "pages": [],
                }

        return {
            "text": native_text,
            "pdf_path": path,
            "source_type": "pdf",
            "method": "native_pymupdf",
            "scanned": False,
            "pages": native_pages,
        }

    def _extract_pdf_native(
        self,
        path: str,
    ) -> Tuple[str, int, List[Dict[str, Any]]]:

        text_parts: List[str] = []
        structured_pages: List[Dict[str, Any]] = []

        with fitz.open(path) as doc:
            page_count = doc.page_count

            for page_index, page in enumerate(doc):
                page_text = self._sanitize_text(
                    page.get_text("text") or ""
                )

                text_parts.append(page_text)

                structured_pages.append(
                    {
                        "page_number": page_index + 1,
                        "text": page_text,
                    }
                )

        combined_text = self._sanitize_text(
            "\n\n".join(p for p in text_parts if p)
        )

        return combined_text, page_count, structured_pages

    # ------------------------------------------------
    # TXT Extraction
    # ------------------------------------------------

    @with_request_id
    def _extract_txt_result(
        self,
        path: str,
        request_id: Optional[str] = None,
    ) -> Dict[str, Any]:

        text = self._extract_txt(path)

        debug_id(
            f"[ContentExtractionService] TXT extraction | chars={len(text)}",
            request_id,
        )

        return {
            "text": text,
            "pdf_path": None,
            "source_type": "txt",
            "method": "native_txt",
            "scanned": False,
            "pages": None,
        }

    def _extract_txt(self, path: str) -> str:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return self._sanitize_text(f.read())

    # ------------------------------------------------
    # CSV Extraction
    # ------------------------------------------------

    @with_request_id
    def _extract_csv(
        self,
        path: str,
        request_id: Optional[str] = None,
    ) -> Optional[Dict[str, Any]]:

        try:
            with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
                reader = csv.reader(f)
                rows = [
                    row
                    for row in reader
                    if any(self._sanitize_text(cell) for cell in row)
                ]

            if not rows:
                warning_id(
                    f"[ContentExtractionService] CSV has no data rows: {path}",
                    request_id,
                )
                return None

            headers = [self._sanitize_text(h) for h in rows[0]]
            chunks: List[str] = []

            for row in rows[1:]:
                pairs = []

                for h, v in zip(headers, row):
                    header = self._sanitize_text(h)
                    value = self._sanitize_text(v)

                    if header and value:
                        pairs.append(f"{header}: {value}")

                if pairs:
                    chunks.append(" | ".join(pairs))

            if not chunks:
                warning_id(
                    f"[ContentExtractionService] CSV produced no text chunks: {path}",
                    request_id,
                )
                return None

            text = self._sanitize_text("\n".join(chunks))

            debug_id(
                f"[ContentExtractionService] CSV extraction | rows={len(chunks)} | chars={len(text)}",
                request_id,
            )

            return {
                "text": text,
                "pdf_path": None,
                "source_type": "csv",
                "method": "native_csv",
                "scanned": False,
                "pages": None,
            }

        except Exception as e:
            error_id(
                f"[ContentExtractionService] CSV extraction failed: {e}",
                request_id,
                exc_info=True,
            )
            return None

    # ------------------------------------------------
    # Excel Extraction
    # ------------------------------------------------

    @with_request_id
    def _extract_excel(
        self,
        path: str,
        request_id: Optional[str] = None,
    ) -> Optional[Dict[str, Any]]:
        """
        Enhanced Excel extraction.

        Handles:
          - multiple sheets
          - normal worksheet cell text
          - merged-cell layouts
          - troubleshooting-style spreadsheets that are not clean tables
          - Excel drawing/textbox text stored in xl/drawings/*.xml
        """

        workbook = None

        try:
            import zipfile
            import xml.etree.ElementTree as ET

            import openpyxl

            workbook = openpyxl.load_workbook(
                filename=path,
                data_only=True,
                read_only=False,
            )

            sheet_names = list(workbook.sheetnames)

            extracted_sections: List[str] = []
            worksheet_line_count = 0
            textbox_line_count = 0

            for worksheet in workbook.worksheets:
                sheet_lines: List[str] = []

                sheet_lines.append("")
                sheet_lines.append("====================")
                sheet_lines.append(f"SHEET: {self._sanitize_text(worksheet.title)}")
                sheet_lines.append("====================")

                for row in worksheet.iter_rows(values_only=True):
                    values = [
                        self._sanitize_text(value)
                        for value in row
                        if value is not None and self._sanitize_text(value)
                    ]

                    if not values:
                        continue

                    line = self._sanitize_text(" | ".join(values))

                    if line:
                        sheet_lines.append(line)
                        worksheet_line_count += 1

                if len(sheet_lines) > 4:
                    extracted_sections.append("\n".join(sheet_lines))

            textbox_sections: List[str] = []

            try:
                if zipfile.is_zipfile(path):
                    with zipfile.ZipFile(path, "r") as archive:
                        drawing_files = [
                            name
                            for name in archive.namelist()
                            if name.startswith("xl/drawings/")
                            and name.endswith(".xml")
                        ]

                        drawing_text_tag = (
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
                        )

                        for drawing_file in drawing_files:
                            try:
                                xml_bytes = archive.read(drawing_file)
                                root = ET.fromstring(xml_bytes)

                                drawing_texts = [
                                    self._sanitize_text(node.text)
                                    for node in root.iter(drawing_text_tag)
                                    if node.text and self._sanitize_text(node.text)
                                ]

                                if not drawing_texts:
                                    continue

                                textbox_sections.append("")
                                textbox_sections.append("====================")
                                textbox_sections.append(
                                    f"TEXTBOX CONTENT: {self._sanitize_text(drawing_file)}"
                                )
                                textbox_sections.append("====================")

                                for text_item in drawing_texts:
                                    cleaned_text_item = self._sanitize_text(text_item)
                                    if cleaned_text_item:
                                        textbox_sections.append(cleaned_text_item)
                                        textbox_line_count += 1

                            except Exception as drawing_error:
                                warning_id(
                                    f"[ContentExtractionService] Failed reading Excel drawing XML "
                                    f"| file={drawing_file} | error={drawing_error}",
                                    request_id,
                                )
                                continue

            except Exception as textbox_error:
                warning_id(
                    f"[ContentExtractionService] Excel textbox extraction skipped: {textbox_error}",
                    request_id,
                )

            if textbox_sections:
                extracted_sections.append("\n".join(textbox_sections))

            text = self._sanitize_text(
                "\n\n".join(
                    section
                    for section in extracted_sections
                    if section and section.strip()
                )
            )

            if not text:
                warning_id(
                    f"[ContentExtractionService] Excel produced no extractable text: {path}",
                    request_id,
                )
                return None

            debug_id(
                f"[ContentExtractionService] Enhanced Excel extraction "
                f"| sheets={len(sheet_names)} "
                f"| worksheet_lines={worksheet_line_count} "
                f"| textbox_lines={textbox_line_count} "
                f"| chars={len(text)}",
                request_id,
            )

            return {
                "text": text,
                "pdf_path": None,
                "source_type": "xlsx",
                "method": "enhanced_excel_cells_and_textboxes",
                "scanned": False,
                "pages": None,
            }

        except Exception as e:
            error_id(
                f"[ContentExtractionService] Excel extraction failed: {e}",
                request_id,
                exc_info=True,
            )
            return None

        finally:
            try:
                if workbook is not None:
                    workbook.close()
            except Exception as close_error:
                warning_id(
                    f"[ContentExtractionService] Failed closing Excel workbook: {close_error}",
                    request_id,
                )

    # ------------------------------------------------
    # DOCX Extraction
    # ------------------------------------------------

    @with_request_id
    def _extract_docx(
        self,
        path: str,
        request_id: Optional[str] = None,
    ) -> Optional[Dict[str, Any]]:

        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(path)
            parts: List[str] = []

            for para in doc.paragraphs:
                txt = self._sanitize_text(para.text)
                if txt:
                    parts.append(txt)

            for table in doc.tables:
                for row in table.rows:
                    cells = [
                        self._sanitize_text(cell.text)
                        for cell in row.cells
                        if self._sanitize_text(cell.text)
                    ]

                    if cells:
                        parts.append(" | ".join(cells))

            text = self._sanitize_text("\n".join(parts))

            if not text:
                warning_id(
                    f"[ContentExtractionService] DOCX produced no extractable text: {path}",
                    request_id,
                )
                return None

            debug_id(
                f"[ContentExtractionService] DOCX extraction | chars={len(text)}",
                request_id,
            )

            return {
                "text": text,
                "pdf_path": None,
                "source_type": "docx",
                "method": "native_docx",
                "scanned": False,
                "pages": None,
            }

        except ImportError:
            warning_id(
                "[ContentExtractionService] python-docx is not installed; DOCX native extraction unavailable",
                request_id,
            )
            return None

        except Exception as e:
            error_id(
                f"[ContentExtractionService] DOCX extraction failed: {e}",
                request_id,
                exc_info=True,
            )
            return None

    # ------------------------------------------------
    # PPTX Extraction
    # ------------------------------------------------

    @with_request_id
    def _extract_pptx(
        self,
        path: str,
        request_id: Optional[str] = None,
    ) -> Optional[Dict[str, Any]]:
        """
        Extract text from PowerPoint .pptx files.

        Handles:
          - slide text boxes
          - slide titles
          - grouped shapes
          - tables
          - speaker notes when available

        Note:
          python-pptx imports as `pptx`.
        """

        try:
            from pptx import Presentation

            presentation = Presentation(path)

            slide_sections: List[str] = []
            slide_count = 0
            text_line_count = 0
            table_line_count = 0
            notes_line_count = 0

            for slide_number, slide in enumerate(presentation.slides, start=1):
                slide_count += 1
                slide_lines: List[str] = []

                slide_lines.append("")
                slide_lines.append("====================")
                slide_lines.append(f"SLIDE: {slide_number}")
                slide_lines.append("====================")

                shape_texts = self._extract_pptx_shapes_text(
                    slide.shapes,
                    request_id=request_id,
                )

                for shape_text in shape_texts:
                    cleaned_shape_text = self._sanitize_text(shape_text)
                    if cleaned_shape_text:
                        slide_lines.append(cleaned_shape_text)
                        text_line_count += 1

                table_texts = self._extract_pptx_tables_text(
                    slide.shapes,
                    request_id=request_id,
                )

                for table_text in table_texts:
                    cleaned_table_text = self._sanitize_text(table_text)
                    if cleaned_table_text:
                        slide_lines.append(cleaned_table_text)
                        table_line_count += 1

                notes_texts = self._extract_pptx_notes_text(
                    slide,
                    request_id=request_id,
                )

                if notes_texts:
                    slide_lines.append("")
                    slide_lines.append("SPEAKER NOTES:")

                    for note_text in notes_texts:
                        cleaned_note_text = self._sanitize_text(note_text)
                        if cleaned_note_text:
                            slide_lines.append(cleaned_note_text)
                            notes_line_count += 1

                if len(slide_lines) > 4:
                    slide_sections.append("\n".join(slide_lines))

            text = self._sanitize_text(
                "\n\n".join(
                    section
                    for section in slide_sections
                    if section and section.strip()
                )
            )

            if not text:
                warning_id(
                    f"[ContentExtractionService] PPTX produced no extractable text: {path}",
                    request_id,
                )
                return None

            debug_id(
                f"[ContentExtractionService] PPTX extraction "
                f"| slides={slide_count} "
                f"| text_lines={text_line_count} "
                f"| table_lines={table_line_count} "
                f"| notes_lines={notes_line_count} "
                f"| chars={len(text)}",
                request_id,
            )

            return {
                "text": text,
                "pdf_path": None,
                "source_type": "pptx",
                "method": "native_pptx",
                "scanned": False,
                "pages": None,
            }

        except ImportError:
            warning_id(
                "[ContentExtractionService] python-pptx is not installed; PPTX native extraction unavailable",
                request_id,
            )
            return None

        except Exception as e:
            error_id(
                f"[ContentExtractionService] PPTX extraction failed: {e}",
                request_id,
                exc_info=True,
            )
            return None

    def _extract_pptx_shapes_text(
        self,
        shapes,
        request_id: Optional[str] = None,
    ) -> List[str]:
        """
        Extract text from normal PowerPoint shapes.

        Also supports grouped shapes by walking nested shape collections.
        """

        extracted: List[str] = []

        for shape in shapes:
            try:
                if hasattr(shape, "shapes"):
                    extracted.extend(
                        self._extract_pptx_shapes_text(
                            shape.shapes,
                            request_id=request_id,
                        )
                    )

                if hasattr(shape, "text"):
                    text = self._sanitize_text(shape.text)
                    if text:
                        extracted.append(text)

            except Exception as shape_error:
                warning_id(
                    f"[ContentExtractionService] Failed reading PPTX shape text: {shape_error}",
                    request_id,
                )
                continue

        return extracted

    def _extract_pptx_tables_text(
        self,
        shapes,
        request_id: Optional[str] = None,
    ) -> List[str]:
        """
        Extract text from PowerPoint tables.

        Also supports grouped shapes by walking nested shape collections.
        """

        extracted: List[str] = []

        for shape in shapes:
            try:
                if hasattr(shape, "shapes"):
                    extracted.extend(
                        self._extract_pptx_tables_text(
                            shape.shapes,
                            request_id=request_id,
                        )
                    )

                if not getattr(shape, "has_table", False):
                    continue

                table = shape.table

                for row in table.rows:
                    values = [
                        self._sanitize_text(cell.text)
                        for cell in row.cells
                        if self._sanitize_text(cell.text)
                    ]

                    if values:
                        extracted.append(" | ".join(values))

            except Exception as table_error:
                warning_id(
                    f"[ContentExtractionService] Failed reading PPTX table text: {table_error}",
                    request_id,
                )
                continue

        return extracted

    def _extract_pptx_notes_text(
        self,
        slide,
        request_id: Optional[str] = None,
    ) -> List[str]:
        """
        Extract speaker notes from a PowerPoint slide when available.
        """

        extracted: List[str] = []

        try:
            notes_slide = getattr(slide, "notes_slide", None)

            if notes_slide is None:
                return extracted

            notes_text_frame = getattr(notes_slide, "notes_text_frame", None)

            if notes_text_frame is not None:
                text = self._sanitize_text(notes_text_frame.text)
                if text:
                    extracted.append(text)

            try:
                for shape in notes_slide.shapes:
                    if hasattr(shape, "text"):
                        text = self._sanitize_text(shape.text)
                        if text and text not in extracted:
                            extracted.append(text)

            except Exception:
                pass

        except Exception as notes_error:
            warning_id(
                f"[ContentExtractionService] Failed reading PPTX speaker notes: {notes_error}",
                request_id,
            )

        return extracted

    # ------------------------------------------------
    # Legacy DOC Handling
    # ------------------------------------------------

    @with_request_id
    def _extract_legacy_doc(
        self,
        path: str,
        request_id: Optional[str] = None,
    ) -> Optional[Dict[str, Any]]:

        warning_id(
            "[ContentExtractionService] Legacy .doc extraction is not supported directly. "
            "Convert .doc to .pdf or .docx in DocumentConversionService first.",
            request_id,
        )
        return None

    # ------------------------------------------------
    # Legacy PPT Handling
    # ------------------------------------------------

    @with_request_id
    def _extract_legacy_ppt(
        self,
        path: str,
        request_id: Optional[str] = None,
    ) -> Optional[Dict[str, Any]]:

        warning_id(
            "[ContentExtractionService] Legacy .ppt extraction is not supported directly. "
            "Convert .ppt to .pptx or .pdf in DocumentConversionService first.",
            request_id,
        )
        return None

    # ------------------------------------------------
    # Scanned Detection Heuristic
    # ------------------------------------------------

    def _is_likely_scanned(self, text: str, pages: int) -> bool:
        if pages <= 0:
            return True

        total_chars = len((text or "").strip())

        if total_chars < 200:
            return True

        avg_chars_per_page = total_chars // max(pages, 1)

        if avg_chars_per_page < 30:
            return True

        return False